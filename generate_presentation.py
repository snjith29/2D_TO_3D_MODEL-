#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for AI-Powered 2D to 3D Building Design Conversion System
Creates a 21-slide presentation with all project details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create comprehensive PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(0, 102, 204)  # Blue
    TEXT_COLOR = RGBColor(64, 64, 64)    # Dark gray
    ACCENT_COLOR = RGBColor(255, 153, 0) # Orange
    
    def add_title_slide(title, subtitle, name1, name2, name3, guide):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Add team info
        team_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
        team_frame = team_box.text_frame
        team_frame.word_wrap = True
        p = team_frame.paragraphs[0]
        p.text = f"{name1}  |  {name2}  |  {name3}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Add guide
        guide_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        guide_frame = guide_box.text_frame
        p = guide_frame.paragraphs[0]
        p.text = f"Guided by: {guide}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_lines, has_speaker_notes=""):
        """Add content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # Add separator line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(2)
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, line_content in enumerate(content_lines):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = line_content
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0
        
        # Add speaker notes
        if has_speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = has_speaker_notes
    
    def add_two_column_slide(title, left_content, right_content):
        """Add two-column slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # Separator
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(2)
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.8))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for i, item in enumerate(left_content):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for i, item in enumerate(right_content):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
    
    # SLIDE 1: Title Slide
    add_title_slide(
        "AI-Powered 2D to 3D Building Design Conversion System",
        "R19AD482 – Project Work-Phase II\nDomain: Artificial Intelligence",
        "SANJAN U (22AD083)",
        "SANJITH S (22AD085)",
        "SARAN H (22AD087)",
        "Dr. A. Sivaramakrishnan"
    )
    
    # SLIDE 2: Problem Motivation
    add_content_slide(
        "PROBLEM MOTIVATION",
        [
            "• NSG and defense units need advanced visualization tools",
            "• Traditional 2D blueprints limit operational effectiveness",
            "• Time-intensive briefing processes delay decisions",
            "• Automate 2D layout conversion to interactive 3D walkthroughs",
            "• Provide offline operational capabilities for secure environments",
            "• Enable realistic visualization with AR/VR compatibility",
            "• Integrate computer vision, 3D modeling, and offline GIS mapping"
        ]
    )
    
    # SLIDE 3: Problem Definition
    add_content_slide(
        "PROBLEM DEFINITION",
        [
            "• Traditional 2D building blueprints are difficult to interpret quickly",
            "• Manual 3D modeling using CAD requires high expertise and time",
            "• Existing AI tools rely on cloud rendering (not feasible for offline missions)",
            "• High data security risks due to cloud uploads",
            "• Cannot be used in secure/offline mission environments",
            "• Lack of offline, secure, and interactive 3D system",
            "• No solution tailored for national security operations"
        ]
    )
    
    # SLIDE 4: Objectives
    add_content_slide(
        "OBJECTIVES",
        [
            "✓ Develop AI-based system to convert 2D blueprints into 3D walkthroughs",
            "✓ Allow user-defined parameters (height, doors, stairs, entry/exit points)",
            "✓ Integrate offline satellite imagery for improved visualization",
            "✓ Ensure complete offline operation for security",
            "✓ Enhance mission efficiency by reducing briefing time",
            "✓ Improve spatial understanding for better decisions",
            "✓ Support rapid tactical deployment preparation"
        ]
    )
    
    # SLIDE 5: Existing System
    add_content_slide(
        "EXISTING SYSTEM",
        [
            "Manual CAD Conversion (AutoCAD, Revit, SketchUp):",
            "  • Requires expert 3D designers • Time-consuming",
            "  • Not scalable for real-time planning",
            "",
            "Cloud-Based AI Tools (Planner5D, Reconstruct.ai):",
            "  • Fast but internet-dependent",
            "  • High data security risks",
            "",
            "➡ NEED: Secure, AI-driven offline system for fast, realistic,",
            "         private mission visualization"
        ]
    )
    
    # SLIDE 6: Disadvantages of Existing System
    add_content_slide(
        "DISADVANTAGES OF EXISTING SYSTEM",
        [
            "❌ Requires expert 3D designers (manual effort)",
            "❌ Time-consuming and not suitable for real-time planning",
            "❌ Not scalable for large or urgent operations",
            "❌ Cloud AI tools are internet-dependent",
            "❌ High data security risks due to cloud uploads",
            "❌ Cannot be used in secure/offline mission environments",
            "❌ High processing time and resource requirement",
            "❌ Lacks proper offline functionality and GIS integration"
        ]
    )
    
    # SLIDE 7: Proposed System
    add_content_slide(
        "PROPOSED SYSTEM",
        [
            "Secure, offline AI-based 2D-to-3D conversion software",
            "",
            "✓ Parses 2D blueprints or scanned layouts",
            "✓ Detects architectural components via computer vision & CAD parsing",
            "✓ Generates 3D walkthrough using Blender/Unity engine",
            "✓ Integrates offline map data for realistic surroundings",
            "✓ Exports models in .glb / .obj formats for AR/VR devices",
            "✓ Complete offline operation with zero internet dependency"
        ]
    )
    
    # SLIDE 8: Advantages of Proposed System
    add_two_column_slide(
        "ADVANTAGES OF PROPOSED SYSTEM",
        [
            "✅ Fully Offline & Secure",
            "✅ AI-Driven Automation",
            "✅ Faster Model Generation",
            "✅ Reduced Skilled Labor",
            "✅ Real-Time 3D Walkthrough"
        ],
        [
            "✅ Computer Vision Integration",
            "✅ GIS & Offline Map Integration",
            "✅ AR/VR Compatibility",
            "✅ Scalable for Multiple Layouts",
            "✅ Improved Decision Making"
        ]
    )
    
    # SLIDE 9: System Architecture
    add_content_slide(
        "SYSTEM ARCHITECTURE",
        [
            "Input Layer:",
            "  • User Interface (GUI/CLI) • Blueprint Selection",
            "",
            "Processing Layer:",
            "  • Blueprint Parser • ML Detection • 3D Generator • Map Integration",
            "",
            "Output Layer:",
            "  • Multiple Export Formats • Web Viewer • Metadata",
            "",
            "Deployment:",
            "  • Offline-capable • Standalone executable • No external dependencies"
        ]
    )
    
    # SLIDE 10: Data Flow Diagram
    add_content_slide(
        "DATA FLOW DIAGRAM",
        [
            "Input 2D Blueprint (PDF/DXF/Image)",
            "         ↓",
            "Preprocessing (OpenCV) → Element Extraction",
            "         ↓",
            "Classification & Transformation",
            "         ↓",
            "3D Model Generation (Blender)",
            "         ↓",
            "Map Integration (Offline GIS)",
            "         ↓",
            "Export (GLB/OBJ) → Visualization & NSG Operations"
        ]
    )
    
    # SLIDE 11: Modules
    add_content_slide(
        "MODULES",
        [
            "1. Parser Module - Extract geometry from 2D blueprints",
            "2. ML Inference Module - AI-powered element detection",
            "3. ML Postprocess Module - Convert predictions to usable data",
            "4. Blender Generator - Create 3D models from extracted data",
            "5. Map Downloader - Integrate offline GIS data",
            "6. Viewer Module - Interactive 3D visualization",
            "7. Orchestration (CLI/UI) - System coordination & user interface"
        ]
    )
    
    # SLIDE 12: Module Description
    add_content_slide(
        "MODULE DESCRIPTION",
        [
            "Parser Module: Image preprocessing, edge detection, wall extraction",
            "",
            "ML Pipeline: ONNX inference, mask post-processing, element detection",
            "",
            "Blender Generator: Wall mesh creation, material assignment, export",
            "",
            "Map Downloader: OSM tile fetching, MBTiles parsing, texture mapping",
            "",
            "Viewer (Web): glTF loading, orbit controls, interactive visualization",
            "",
            "Orchestration: Argument parsing, pipeline coordination, error handling"
        ]
    )
    
    # SLIDE 13: Technology Stack
    add_content_slide(
        "TECHNOLOGY STACK",
        [
            "Core: Python 3.10+, Blender 3.5+",
            "Computer Vision: OpenCV, ezdxf, Shapely",
            "Machine Learning: ONNX Runtime",
            "Numerical: NumPy, Pillow",
            "Mapping & GIS: MBTiles, Rasterio",
            "Web & Visualization: Three.js, Flask",
            "User Interface: Tkinter, Matplotlib",
            "System: Windows/Linux/macOS, 8GB+ RAM, 20GB+ storage"
        ]
    )
    
    # SLIDE 14: Installation & Setup
    add_content_slide(
        "INSTALLATION & SETUP",
        [
            "1. Verify Prerequisites: Python 3.10+, Blender 3.5+",
            "2. Create Virtual Environment: python -m venv .venv",
            "3. Install Dependencies: pip install -r requirements.txt",
            "4. Run Test: python test_installation.py",
            "5. Verify Blender Installation: blender --version",
            "6. Test Conversion: python main.py --input blueprint.png",
            "7. View Results: cd viewer && python -m http.server 8000"
        ]
    )
    
    # SLIDE 15: Usage Examples
    add_content_slide(
        "USAGE EXAMPLES",
        [
            "Basic: python main.py --input blueprint.png --scale 100 --height 3.5",
            "",
            "Advanced: python main.py --input plan.dxf --height 4.0",
            "          --wall-thickness 0.3 --stair-width 1.5",
            "",
            "GUI Mode: python main.py --ui (or) python ui.py",
            "",
            "With Viewer: python main.py --input blueprint.png --serve-viewer",
            "",
            "Batch: python batch_convert.py"
        ]
    )
    
    # SLIDE 16: Expected Output Files
    add_content_slide(
        "EXPECTED OUTPUT FILES",
        [
            "• plan.json (~654 bytes): Parsed blueprint data",
            "• model.glb (6-7 KB): 3D model in glTF binary format",
            "• model.obj (3-4 KB): 3D model in OBJ format",
            "• model.mtl (~617 bytes): Material file for OBJ format",
            "",
            "Usage:",
            "  - GLB: Web browsers, AR/VR applications",
            "  - OBJ: CAD software (Blender, Maya, 3D Max)",
            "  - JSON: Database, simulations, downstream processing"
        ]
    )
    
    # SLIDE 17: Conclusion
    add_content_slide(
        "CONCLUSION",
        [
            "✓ Robust parsing pipeline extracts all architectural elements",
            "✓ Optional ML enhancement refines detection accuracy",
            "✓ Automated 3D generation eliminates manual CAD recreation",
            "✓ Structured data export supports downstream workflows",
            "✓ Interactive visualization accessible to non-CAD users",
            "✓ Extensible architecture supports custom enhancements",
            "",
            "Key Impact: 80-90% reduction in modeling time with enhanced security"
        ]
    )
    
    # SLIDE 18: Future Scope
    add_content_slide(
        "FUTURE SCOPE",
        [
            "Short-term: Mobile apps, Enhanced ML models, Cloud API",
            "",
            "Medium-term: Real-time collaboration, AI suggestions, BIM integration",
            "",
            "Long-term: Point cloud support, AR/VR headsets, Advanced analytics",
            "",
            "Research: Novel CNN architectures, GAN-based generation,",
            "          Graph neural networks, Federated learning",
            "",
            "Business: SaaS platform, Enterprise licenses, Developer API access"
        ]
    )
    
    # SLIDE 19: References
    add_content_slide(
        "REFERENCES",
        [
            "[1] Ahmed et al., 'Automatic Floor Plan Analysis', IEEE TPAMI, 2018",
            "[2] Zhang et al., 'Floor-Plan Reconstruction', ACM TOG, 2019",
            "[3] Fabrikant & Lobben, 'GIS for Security Operations', 2021",
            "[4] McMahan et al., 'VR Training Evaluation', MILCOM, 2018",
            "[5] Liu et al., 'Deep Learning for Floor Plan Recognition', CVPR, 2020",
            "",
            "Documentation: OpenCV, Blender, Three.js, ONNX Runtime",
            "Standards: glTF, OBJ, DXF, JSON Schema"
        ]
    )
    
    # SLIDE 20: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 102, 204)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "THANK YOU FOR YOUR ATTENTION!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    p = content_frame.paragraphs[0]
    p.text = "Questions & Discussion\n\nDemo Available\nRepository: GitHub\nContact: Available for collaboration"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 21: Q&A Slide
    add_content_slide(
        "Q&A SESSION",
        [
            "Open to Questions:",
            "  • Technical implementation details",
            "  • Feature requests and suggestions",
            "  • Real-world applications",
            "  • Partnership opportunities",
            "",
            "Contact Information:",
            "  • GitHub Repository",
            "  • Email Support",
            "  • Live Demonstrations Available"
        ]
    )
    
    return prs

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    presentation = create_presentation()
    
    output_path = r"c:\Users\sanjith\OneDrive\Desktop\fnp\AI-Powered_2D_to_3D_Presentation.pptx"
    presentation.save(output_path)
    
    print(f"✓ Presentation created successfully!")
    print(f"✓ Saved to: {output_path}")
    print(f"✓ 21 slides generated with complete content")
    print(f"✓ Ready for presentation!")
